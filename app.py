from flask import Flask, render_template, request, redirect, url_for, session, flash, jsonify, send_file
from werkzeug.security import check_password_hash ,generate_password_hash
from deep_translator import GoogleTranslator
from gtts import gTTS
import os
import yt_dlp
try:
    import whisper
    WHISPER_AVAILABLE = True
except (ImportError, TypeError) as e:
    # Handle Windows libc loading issue or missing whisper
    print(f"Warning: Could not import whisper: {e}")
    WHISPER_AVAILABLE = False
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
import re
from openai import OpenAI
import sqlite3
from werkzeug.utils import secure_filename
from io import BytesIO
import tempfile

# progress tracking helper
def log_progress(user_id, description):
    if not user_id:
        return
    conn = sqlite3.connect("database.db")
    c = conn.cursor()
    c.execute("INSERT INTO progress (user_id, description) VALUES (?, ?)", (user_id, description))
    conn.commit()
    conn.close()


def init_db():
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    c.execute("""
        CREATE TABLE IF NOT EXISTS videos (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            video_id TEXT UNIQUE,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)

    c.execute("""
        CREATE TABLE IF NOT EXISTS notes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            video_id TEXT,
            timestamp TEXT,
            content TEXT
        )
    """)
    c.execute("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL
        )
    """)

    c.execute("""
        CREATE TABLE IF NOT EXISTS history (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            video_id TEXT,
            notes TEXT
        )
    """)
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        tool TEXT,
        title TEXT,
        content TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    # progress logs for activities
    c.execute("""
    CREATE TABLE IF NOT EXISTS progress (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        description TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    #   Check existing columns
    c.execute("PRAGMA table_info(history)")
    columns = [column[1] for column in c.fetchall()]

    if "tool" not in columns:
        c.execute("ALTER TABLE history ADD COLUMN tool TEXT")

    if "title" not in columns:
        c.execute("ALTER TABLE history ADD COLUMN title TEXT")

    if "content" not in columns:
        c.execute("ALTER TABLE history ADD COLUMN content TEXT")

    if "created_at" not in columns:
        c.execute("ALTER TABLE history ADD COLUMN created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP")
    # ensure progress table exists (if added later)
    c.execute("""
    CREATE TABLE IF NOT EXISTS progress (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        description TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )""")
    conn.commit()
    conn.close()

init_db()

client = OpenAI(api_key="YOUR_OPENAI_API_KEY")


app = Flask(__name__)
app.secret_key = "supersecretkey"

def extract_video_id(url):
    pattern = r"(?:v=|\/)([0-9A-Za-z_-]{11}).*"
    match = re.search(pattern, url)
    return match.group(1) if match else None

@app.route('/')
def home():
    print("SESSION DATA:", dict(session))
    return render_template("index.html")
    #return render_template('login.html')

@app.route('/translator', methods=['GET', 'POST'])
def translator():
    translated_text = ""
    original_text = ""
    audio_file = None
    error_message = None

    if request.method == 'POST':
        try:
            original_text = request.form.get('text', '').strip()
            target_language = request.form.get('language', '')
            
            # Validation
            if not original_text:
                error_message = "Please enter text to translate"
                if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return jsonify({
                        'success': False,
                        'error': error_message
                    }), 400
            
            if not target_language:
                error_message = "Please select a language"
                if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return jsonify({
                        'success': False,
                        'error': error_message
                    }), 400
            
            # Translate text
            translated_text = GoogleTranslator(source='auto', target=target_language).translate(original_text)
            
            if not translated_text:
                raise Exception("Translation returned empty result")

            # save translation to history
            user_id = session.get('user_id')
            if user_id:
                conn = sqlite3.connect("database.db")
                c = conn.cursor()
                title = f"Translator to {target_language}"
                c.execute("INSERT INTO history (user_id, tool, title, content) VALUES (?, ?, ?, ?)",
                          (user_id, "Translator", title, translated_text))
                conn.commit()
                conn.close()
                log_progress(user_id, f"Translated text to {target_language}")

            # Generate voice file
            try:
                tts = gTTS(text=translated_text, lang=target_language, slow=False)
                audio_path = os.path.join("static", "output.mp3")
                tts.save(audio_path)
                audio_file = "output.mp3"
            except Exception as audio_error:
                print(f"Audio generation error: {audio_error}")
                # Don't fail translation if audio generation fails
                audio_file = None
            
            # Check if this is an AJAX request
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({
                    'success': True,
                    'translated_text': translated_text,
                    'audio_file': audio_file,
                    'original_text': original_text
                })
        
        except Exception as e:
            error_message = f"Translation error: {str(e)}"
            print(f"Error in translator: {error_message}")
            
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({
                    'success': False,
                    'error': error_message
                }), 500

    return render_template('translator.html',
                           translated_text=translated_text,
                           original_text=original_text,
                           audio_file=audio_file,
                           error_message=error_message)

@app.route("/youtube-notes", methods=["GET", "POST"])
def youtube_notes():
    notes = None
    video_id = None
    mode = None
    ai_mode = None

    if request.method == "POST":
        url = request.form.get("url")
        video_id = extract_video_id(url)

        # Download audio using yt_dlp
        ydl_opts = {
            'format': 'bestaudio/best',
            'outtmpl': 'audio.%(ext)s',
            'quiet': True
        }

        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])

        # Try to locate the downloaded audio file (yt_dlp may use different extensions)
        import glob
        downloaded = glob.glob(os.path.join(os.getcwd(), "audio.*"))
        if not downloaded:
            error_msg = "Downloaded audio file not found (yt_dlp output didn't match audio.*)."
            print(error_msg)
            notes = ""
            history = []
            if "user_id" in session:
                conn = sqlite3.connect("database.db")
                c = conn.cursor()
                c.execute("""
                SELECT tool, title, created_at
                FROM history
                WHERE user_id = ?
                ORDER BY created_at DESC
                LIMIT 5
                """, (session["user_id"],))
                history = c.fetchall()
                conn.close()
            return render_template("youtube.html", notes=notes, video_id=video_id, history=history, mode=mode, ai_mode=ai_mode)

        audio_path = downloaded[0]

        # Check if whisper is available
        if not WHISPER_AVAILABLE:
            error_msg = "Whisper audio transcription is not available on this system. Please check the installation."
            print(error_msg)
            notes = ""
            history = []
            if "user_id" in session:
                conn = sqlite3.connect("database.db")
                c = conn.cursor()
                c.execute("""
                SELECT tool, title, created_at
                FROM history
                WHERE user_id = ?
                ORDER BY created_at DESC
                LIMIT 5
                """, (session["user_id"],))
                history = c.fetchall()
                conn.close()
            return render_template("youtube.html", notes=notes, video_id=video_id, history=history, mode=mode, ai_mode=ai_mode)

        # Load whisper model and transcribe
        model = whisper.load_model("base")
        try:
            result = model.transcribe(audio_path)
        finally:
            # cleanup downloaded audio if present
            try:
                os.remove(audio_path)
            except Exception:
                pass
        mode = request.form.get("mode")
        ai_mode = request.form.get("ai_mode", "short")

        if mode == "general":
            formatted_notes = ""
            conn = sqlite3.connect("database.db")
            c = conn.cursor()

            # Insert notes line by line
            for segment in result["segments"]:
                start = int(segment["start"])
                minutes = start // 60
                seconds = start % 60
                timestamp = f"{minutes:02d}:{seconds:02d}"

                text = segment["text"].strip()

                c.execute(
                "INSERT INTO notes (video_id, timestamp, content) VALUES (?, ?, ?)",
                (video_id, timestamp, text)
                )

            conn.commit()
            conn.close()

            for segment in result["segments"]:
                start = int(segment["start"])
                minutes = start // 60
                seconds = start % 60
                timestamp = f"[{minutes:02d}:{seconds:02d}]"

                text = segment["text"].strip()
                formatted_notes += f"{timestamp} {text}\n\n"

            notes = formatted_notes

            # SAVE TO HISTORY
            user_id = session.get("user_id")
            if user_id:
                conn = sqlite3.connect("database.db")
                c = conn.cursor()
                c.execute("""
                    INSERT INTO history (user_id, tool, title, content)
                    VALUES (?, ?, ?, ?)
                     """, (user_id, "YouTube Notes", video_id, notes))
                conn.commit()
                conn.close()
                log_progress(user_id, f"Generated normal notes for {video_id}")

        elif mode == "ai":
            # generate AI notes according to ai_mode
            transcript = result.get('text','')
            if not transcript:
                notes = ""  # safety
            else:
                if ai_mode == 'short':
                    prompt = f"Summarize the following text in a concise paragraph:\n\n{transcript}"
                elif ai_mode == 'exam':
                    prompt = f"Create detailed exam-focused revision notes from the following text:\n\n{transcript}"
                elif ai_mode == 'bullet':
                    prompt = f"Convert the following content into bullet-point smart notes:\n\n{transcript}"
                elif ai_mode == 'simple':
                    prompt = f"Explain the following content in simplified terms for a beginner:\n\n{transcript}"
                else:
                    prompt = f"Summarize the following text:\n\n{transcript}"

                try:
                    response = client.chat.completions.create(
                        model="gpt-4o-mini",
                        messages=[{"role": "user", "content": prompt}]
                    )
                    notes = response.choices[0].message.content
                except Exception as e:
                    print("AI summarization error", e)
                    notes = "[Error generating AI notes]"

            # SAVE AI NOTES TO HISTORY
            user_id = session.get("user_id")
            if user_id:
                conn = sqlite3.connect("database.db")
                c = conn.cursor()
                c.execute("""
                    INSERT INTO history (user_id, tool, title, content)
                    VALUES (?, ?, ?, ?)
                """, (user_id, f"AI Notes ({ai_mode})", video_id, notes))
                conn.commit()
                conn.close()
                log_progress(user_id, f"Generated AI notes ({ai_mode}) for {video_id}")

        # -------- AFTER POST LOGIC --------
    #.......FETCH HERE>>>>>>>
    history = []

    if "user_id" in session:
        conn = sqlite3.connect("database.db")
        c = conn.cursor()
        c.execute("""
        SELECT tool, title, created_at
        FROM history
        WHERE user_id = ?
        ORDER BY created_at DESC
        LIMIT 5
        """, (session["user_id"],))
        history = c.fetchall()
        conn.close()

    return render_template("youtube.html",
                           notes=notes,
                           video_id=video_id,
                           history=history,
                           mode=mode,
                           ai_mode=ai_mode)
        
    

@app.route("/history")
def history():
    import sqlite3
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    c.execute("SELECT video_id, created_at FROM videos ORDER BY created_at DESC")
    videos = c.fetchall()

    conn.close()

    return render_template("history.html", videos=videos)

@app.route("/load-video/<video_id>")
def load_video(video_id):
    import sqlite3
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    c.execute("SELECT timestamp, content FROM notes WHERE video_id = ?", (video_id,))
    rows = c.fetchall()

    conn.close()

    formatted_notes = ""
    for row in rows:
        formatted_notes += f"[{row[0]}] {row[1]}\n\n"

    conn = sqlite3.connect("database.db")
    c = conn.cursor()
    c.execute("SELECT video_id FROM videos ORDER BY created_at DESC")
    videos = c.fetchall()
    conn.close()

    return render_template("youtube.html",
                       notes=formatted_notes,
                       video_id=video_id,
                       history=videos)




@app.route("/download_pdf", methods=["POST"])
def download_pdf():
    notes = request.form.get("notes")
    print("NOTES RECEIVED:", notes)   # 👈 ADD THIS

    if not notes:
        return "No notes received!"

    from io import BytesIO
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer)
    elements = []
    styles = getSampleStyleSheet()

    for line in notes.split("\n"):
        stripped = line.strip()
        if stripped.startswith("-") or stripped.startswith("•") or stripped.startswith("*"):
            elements.append(Paragraph(stripped, styles.get("Bullet", styles["Normal"])))
        else:
            elements.append(Paragraph(line, styles["Normal"]))
        elements.append(Spacer(1, 0.2 * inch))

    doc.build(elements)
    buffer.seek(0)

    # progress log
    user_id = session.get('user_id')
    log_progress(user_id, 'Downloaded notes PDF')

    return send_file(
        buffer,
        as_attachment=True,
        download_name="Video_Notes.pdf",
        mimetype="application/pdf"
    )

@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        username = request.form["username"]
        password = generate_password_hash(request.form["password"])

        conn = sqlite3.connect("database.db")
        c = conn.cursor()

        try:
            c.execute("INSERT INTO users (username, password) VALUES (?, ?)", (username, password))
            conn.commit()
            flash("Registration successful! Please login.")
            return redirect(url_for("login"))
        except:
            flash("Username already exists.")

        conn.close()

    return render_template("register.html")



@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":

        username = request.form.get("username").strip()
        password = request.form.get("password").strip()

        conn = sqlite3.connect("database.db")
        c = conn.cursor()

        c.execute("SELECT * FROM users WHERE username=?", (username,))
        user = c.fetchone()

        conn.close()

        if user is None:
            return render_template("login.html", error="User not found")

        stored_password = user[2]

        # ✅ THIS IS THE FIX
        if check_password_hash(stored_password, password):
            session["user_id"] = user[0]
            session["username"] = user[1]
            return redirect("/")
        else:
            return render_template("login.html", error="Wrong password")

    return render_template("login.html")

@app.route("/profile")
def profile():
    # profile page simply redirects to dashboard for now
    if "user_id" not in session:
        return redirect(url_for("login"))
    return redirect(url_for("dashboard"))


@app.route("/dashboard")
def dashboard():
    if "user_id" not in session:
        return redirect(url_for("login"))
    conn = sqlite3.connect("database.db")
    c = conn.cursor()
    c.execute("SELECT tool, title, created_at FROM history WHERE user_id = ? ORDER BY created_at DESC", (session["user_id"],))
    history = c.fetchall()
    c.execute("SELECT description, created_at FROM progress WHERE user_id = ? ORDER BY created_at DESC", (session["user_id"],))
    progress_items = c.fetchall()
    conn.close()
    return render_template("dashboard.html", history=history, progress=progress_items)


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))


@app.route('/team')
def team():
    # simple static page describing team members
    return render_template('team.html')


@app.route('/pdf-tools')
def pdf_tools():
    # legacy route; redirect to first converter page
    return redirect(url_for('pdf_to_word_page'))

@app.route('/pdf-to-word')
def pdf_to_word_page():
    return render_template('pdf_to_word.html')

@app.route('/word-to-pdf')
def word_to_pdf_page():
    return render_template('word_to_pdf.html')

@app.route('/ppt-to-pdf')
def ppt_to_pdf_page():
    return render_template('ppt_to_pdf.html')

@app.route('/resume-builder', methods=['GET','POST'])
def resume_builder():
    if request.method == 'POST':
        name = request.form.get('name','').strip()
        email = request.form.get('email','').strip()
        phone = request.form.get('phone','').strip()
        education = request.form.get('education','').strip()
        skills = request.form.get('skills','').strip()
        experience = request.form.get('experience','').strip()

        user_id = session.get('user_id')
        log_progress(user_id, 'Generated Resume')

        # Generate professional PDF resume
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.platypus import Table, TableStyle, PageBreak
        from reportlab.lib.colors import HexColor

        buffer = BytesIO()
        page_width, page_height = letter
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=0.6*inch,
            leftMargin=0.6*inch,
            topMargin=0.5*inch,
            bottomMargin=0.5*inch
        )

        styles = getSampleStyleSheet()
        elements = []

        # Create custom styles
        title_style = styles['Heading1']
        title_style.alignment = 1  # CENTER
        title_style.fontSize = 24
        title_style.textColor = HexColor('#2c2c2c')
        title_style.spaceAfter = 6

        section_style = styles.add(ParagraphStyle(
            name='SectionTitle',
            parent=styles['Heading2'],
            fontSize=12,
            textColor=HexColor('#667eea'),
            spaceAfter=6,
            spaceBefore=10,
            fontName='Helvetica-Bold',
            borderPadding=3,
            borderColor=HexColor('#667eea'),
            borderWidth=0,
            borderRadius=2
        ))

        contact_style = styles.add(ParagraphStyle(
            name='ContactInfo',
            parent=styles['Normal'],
            fontSize=9.5,
            textColor=HexColor('#666666'),
            alignment=1,  # CENTER
            spaceAfter=12
        ))

        content_style = styles.add(ParagraphStyle(
            name='Content',
            parent=styles['Normal'],
            fontSize=10,
            textColor=HexColor('#444444'),
            spaceAfter=6,
            leading=13
        ))

        from reportlab.pdfgen import canvas
        from reportlab.lib.styles import ParagraphStyle

        # Header with name
        elements.append(Paragraph(f"<b>{name.upper()}</b>", title_style))

        # Contact information
        contact_info = []
        if email:
            contact_info.append(email)
        if phone:
            contact_info.append(phone)
        
        if contact_info:
            elements.append(Paragraph(" | ".join(contact_info), contact_style))

        elements.append(Spacer(1, 0.15*inch))

        # Education Section
        if education:
            elements.append(Paragraph("EDUCATION", section_style))
            for line in education.split('\n'):
                if line.strip():
                    elements.append(Paragraph(line.strip(), content_style))
            elements.append(Spacer(1, 0.08*inch))

        # Skills Section
        if skills:
            elements.append(Paragraph("SKILLS", section_style))
            skill_items = []
            for line in skills.split('\n'):
                if line.strip():
                    skill_items.append(Paragraph(f"• {line.strip()}", content_style))
            elements.append(Spacer(1, 0.02*inch))
            for item in skill_items:
                elements.append(item)
            elements.append(Spacer(1, 0.08*inch))

        # Experience Section
        if experience:
            elements.append(Paragraph("EXPERIENCE & PROJECTS", section_style))
            for line in experience.split('\n'):
                if line.strip():
                    elements.append(Paragraph(line.strip(), content_style))

        # Build PDF
        doc.build(elements)
        buffer.seek(0)

        return send_file(
            buffer,
            as_attachment=True,
            download_name=f'{name.replace(" ", "_")}_resume.pdf',
            mimetype='application/pdf'
        )

    return render_template('resume_builder.html')


@app.route('/generate-resume', methods=['POST'])
def generate_resume():
    """Generate professional PDF resume with photo support"""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.platypus import Table, TableStyle, PageBreak, Image
    from reportlab.lib.colors import HexColor
    from PIL import Image as PILImage
    import base64

    try:
        name = request.form.get('name', '').strip()
        email = request.form.get('email', '').strip()
        phone = request.form.get('phone', '').strip()
        location = request.form.get('location', '').strip()
        summary = request.form.get('summary', '').strip()
        degree = request.form.get('degree', '').strip()
        university = request.form.get('university', '').strip()
        graduation = request.form.get('graduation', '').strip()
        skills = request.form.get('skills', '').strip()
        experience = request.form.get('experience', '').strip()
        certifications = request.form.get('certifications', '').strip()
        languages = request.form.get('languages', '').strip()
        photo_data = request.form.get('photo', '').strip()

        user_id = session.get('user_id')
        log_progress(user_id, 'Generated Professional Resume')

        buffer = BytesIO()
        page_width, page_height = letter
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=0.5*inch,
            leftMargin=0.5*inch,
            topMargin=0.6*inch,
            bottomMargin=0.5*inch
        )

        styles = getSampleStyleSheet()
        elements = []

        # Custom styles
        title_style = ParagraphStyle(
            name='Title',
            parent=styles['Heading1'],
            fontSize=26,
            textColor=HexColor('#2c2c2c'),
            spaceAfter=4,
            spaceBefore=0,
            fontName='Helvetica-Bold',
            alignment=0
        )

        section_style = ParagraphStyle(
            name='SectionTitle',
            parent=styles['Heading2'],
            fontSize=11,
            textColor=HexColor('#667eea'),
            spaceAfter=8,
            spaceBefore=12,
            fontName='Helvetica-Bold',
            textTransform='uppercase',
            letterSpacing=2
        )

        contact_style = ParagraphStyle(
            name='Contact',
            parent=styles['Normal'],
            fontSize=9,
            textColor=HexColor('#666666'),
            spaceAfter=2,
            leading=11
        )

        content_style = ParagraphStyle(
            name='Content',
            parent=styles['Normal'],
            fontSize=9.5,
            textColor=HexColor('#555555'),
            spaceAfter=4,
            leading=12,
            leftIndent=12
        )

        content_title_style = ParagraphStyle(
            name='ContentTitle',
            parent=styles['Normal'],
            fontSize=10,
            textColor=HexColor('#2c2c2c'),
            fontName='Helvetica-Bold',
            spaceAfter=2,
            leading=12
        )

        # Build header with photo
        header_data = []
        header_row = []

        # Add photo if available
        if photo_data and photo_data.startswith('data:image'):
            try:
                # Extract base64 data
                header_b64 = photo_data.split(',')[1]
                photo_bytes = base64.b64decode(header_b64)
                photo_buffer = BytesIO(photo_bytes)
                
                # Resize photo
                with PILImage.open(photo_buffer) as img:
                    img_resized = img.resize((100, 100), PILImage.Resampling.LANCZOS)
                    photo_buffer_resized = BytesIO()
                    img_resized.save(photo_buffer_resized, format='PNG')
                    photo_buffer_resized.seek(0)
                    
                    photo_img = Image(photo_buffer_resized, width=0.9*inch, height=0.9*inch)
                    header_row.append(photo_img)
            except Exception as e:
                print(f"Error processing photo: {e}")

        # Create name and contact info cell without emojis
        contact_text = f"<b>{name.upper()}</b>"
        contact_parts = []
        if location:
            contact_parts.append(location)
        if phone:
            contact_parts.append(phone)
        if email:
            contact_parts.append(email)
        if contact_parts:
            contact_text += "<br/>" + " | ".join(contact_parts)
        contact_para = Paragraph(contact_text, contact_style)
        header_row.append(contact_para)

        # Build header (with photo if present)
        if len(header_row) == 2:
            # photo + contact
            header_table = Table([header_row], colWidths=[1.0*inch, 4.5*inch])
            header_table.setStyle(TableStyle([
                ('LEFTPADDING', (0, 0), (-1, -1), 0),
                ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                ('TOPPADDING', (0, 0), (-1, -1), 0),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 0),
                ('VALIGN', (0, 0), (0, 0), 'TOP'),
                ('VALIGN', (1, 0), (1, 0), 'TOP'),
            ]))
            elements.append(header_table)
        elif len(header_row) == 1:
            # only contact info
            elements.append(header_row[0])
        else:
            elements.append(Paragraph(f"<b>{name.upper()}</b>", title_style))
            if location or phone or email:
                contact_parts = []
                if location:
                    contact_parts.append(location)
                if phone:
                    contact_parts.append(phone)
                if email:
                    contact_parts.append(email)
                elements.append(Paragraph(" | ".join(contact_parts), contact_style))

        elements.append(Spacer(1, 0.15*inch))

        # Professional Summary
        if summary:
            elements.append(Paragraph("PROFESSIONAL SUMMARY", section_style))
            elements.append(Paragraph(summary, content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Education
        if degree or university or graduation:
            elements.append(Paragraph("EDUCATION", section_style))
            edu_text = f"<b>{degree}</b>"
            if university or graduation:
                edu_text += f"<br/>{university}"
                if graduation:
                    edu_text += f" ({graduation})"
            elements.append(Paragraph(edu_text, content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Skills
        if skills:
            elements.append(Paragraph("SKILLS", section_style))
            skills_list = [s.strip() for s in skills.split(',') if s.strip()]
            for i, skill in enumerate(skills_list):
                if i < len(skills_list) - 1:
                    elements.append(Paragraph(f"{skill} •", content_style))
                else:
                    elements.append(Paragraph(skill, content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Experience & Projects
        if experience:
            elements.append(Paragraph("EXPERIENCE & PROJECTS", section_style))
            sections = experience.split('\n\n')
            for section in sections:
                lines = section.strip().split('\n')
                if lines:
                    elements.append(Paragraph(f"<b>{lines[0]}</b>", content_title_style))
                    for line in lines[1:]:
                        if line.strip():
                            elements.append(Paragraph(f"• {line.strip()}", content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Certifications
        if certifications:
            elements.append(Paragraph("CERTIFICATIONS", section_style))
            for cert in certifications.split('\n'):
                if cert.strip():
                    elements.append(Paragraph(f"• {cert.strip()}", content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Languages
        if languages:
            elements.append(Paragraph("LANGUAGES", section_style))
            for lang in languages.split('\n'):
                if lang.strip():
                    elements.append(Paragraph(f"• {lang.strip()}", content_style))

        # Build PDF
        doc.build(elements)
        buffer.seek(0)

        return send_file(
            buffer,
            as_attachment=True,
            download_name=f'{name.replace(" ", "_")}_resume.pdf',
            mimetype='application/pdf'
        )

    except Exception as e:
        print(f"Error generating resume: {e}")
        return jsonify({'error': str(e)}), 500


# ------ AI study assistant endpoints ------
@app.route('/ai/quiz', methods=['POST'])
def ai_quiz():
    data = request.get_json() or {}
    notes = data.get('notes', '')
    try:
        prompt = (
            "Generate a set of study questions based on the following notes. "
            "Provide 5 multiple-choice questions (with four options each, indicate correct answer and explanation), "
            "5 true/false questions with correct answers, and 3 long-answer questions with answers. "
            f"Return JSON with keys mcq, tf, long.\n\nNotes:\n{notes}"
        )
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        text = resp.choices[0].message.content
        user_id = session.get('user_id')
        log_progress(user_id, 'Generated AI quiz')
        return jsonify({'quiz': text})
    except Exception as e:
        print("AI quiz error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/flashcards', methods=['POST'])
def ai_flashcards():
    data = request.get_json() or {}
    notes = data.get('notes', '')
    try:
        prompt = (
            "Create a list of concise flashcards (question and answer pairs) from the following notes. "
            f"Notes:\n{notes}"
        )
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        user_id = session.get('user_id')
        log_progress(user_id, 'Created AI flashcards')
        return jsonify({'flashcards': resp.choices[0].message.content})
    except Exception as e:
        print("AI flashcards error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/mindmap', methods=['POST'])
def ai_mindmap():
    data = request.get_json() or {}
    notes = data.get('notes', '')
    try:
        prompt = (
            "Generate a structured mind-map description (use indentation or bullet levels) "
            "based on these notes.\n" + notes
        )
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        user_id = session.get('user_id')
        log_progress(user_id, 'Generated AI mindmap')
        return jsonify({'mindmap': resp.choices[0].message.content})
    except Exception as e:
        print("AI mindmap error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/doubt', methods=['POST'])
def ai_doubt():
    data = request.get_json() or {}
    notes = data.get('notes', '')
    question = data.get('question', '')
    try:
        prompt = (
            "You are an intelligent assistant limited to the following notes. "
            "Answer the user question based only on that content. "
            f"Notes:\n{notes}\nQuestion:\n{question}"
        )
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        user_id = session.get('user_id')
        log_progress(user_id, 'Answered AI doubt')
        return jsonify({'answer': resp.choices[0].message.content})
    except Exception as e:
        print("AI doubt error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/study-plan', methods=['POST'])
def ai_study_plan():
    data = request.get_json() or {}
    exam_date = data.get('exam_date', '')
    hours_per_day = data.get('hours_per_day', '')
    try:
        prompt = (
            "Create a personalized study schedule plan based on the exam date "
            f"({exam_date}) and available study hours per day ({hours_per_day}). "
            "Output a day-by-day plan with suggestions."
        )
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        user_id = session.get('user_id')
        log_progress(user_id, 'Generated AI study plan')
        return jsonify({'plan': resp.choices[0].message.content})
    except Exception as e:
        print("AI study plan error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/pdf-summary', methods=['POST'])
def ai_pdf_summary():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    file = request.files['file']
    try:
        import fitz  # pymupdf
        pdf_bytes = file.read()
        doc = fitz.open(stream=pdf_bytes, filetype='pdf')
        full_text = ''
        for page in doc:
            full_text += page.get_text()
        doc.close()

        prompt = (
            "Summarize the following PDF contents into concise bullet points:\n\n" + full_text
        )
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        summary = resp.choices[0].message.content
        user_id = session.get('user_id')
        log_progress(user_id, 'Analyzed PDF via AI')
        return jsonify({'summary': summary})
    except Exception as e:
        print("AI PDF summary error", e)
        return jsonify({'error': str(e)}), 500
@app.route('/convert-pdf-to-word', methods=['POST'])
def convert_pdf_to_word():
    """Convert PDF to Word (.docx) format"""
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({'error': 'Invalid file type'}), 400
    
    try:
        # Import pdf2docx and use Converter class (api changed)
        from pdf2docx import Converter
        
        # Create temporary files
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_file:
            pdf_path = pdf_file.name
            file.save(pdf_path)
        
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as docx_file:
            docx_path = docx_file.name
        
        # Convert PDF to DOCX using new api
        cv = Converter(pdf_path)
        cv.convert(docx_path)
        cv.close()
        
        # Read the converted file
        with open(docx_path, 'rb') as f:
            file_data = f.read()
        
        # Clean up temporary files
        os.remove(pdf_path)
        os.remove(docx_path)
        
        # Return the file
        # record progress
        user_id = session.get('user_id')
        log_progress(user_id, 'Converted PDF to Word')
        return send_file(
            BytesIO(file_data),
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name=secure_filename(file.filename.replace('.pdf', '.docx'))
        )
    
    except ImportError:
        return jsonify({'error': 'pdf2docx library not installed or API incompatible. Please install/update it.'}), 500
    except Exception as e:
        print(f"Error converting PDF to Word: {str(e)}")
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500


@app.route('/convert-word-to-pdf', methods=['POST'])
def convert_word_to_pdf():
    """Convert Word (.docx) to PDF format"""
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '' or not (file.filename.endswith('.docx') or file.filename.endswith('.doc')):
        return jsonify({'error': 'Invalid file type'}), 400
    
    try:
        from docx import Document
        
        # Create temporary files
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as docx_file:
            docx_path = docx_file.name
            file.save(docx_path)
        
        # Load Word document
        doc = Document(docx_path)
        
        # Create PDF
        pdf_buffer = BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        elements = []
        
        # Add paragraphs from Word document
        for para in doc.paragraphs:
            if para.text.strip():
                try:
                    # Handle different paragraph styles
                    style = styles['Normal']
                    if para.style and para.style.name.startswith('Heading'):
                        style = styles['Heading1']
                    
                    elements.append(Paragraph(para.text, style))
                except:
                    elements.append(Paragraph(para.text, styles['Normal']))
                elements.append(Spacer(1, 0.2))
        
        # Add tables if they exist
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_data.append(row_data)
            
            if table_data:
                from reportlab.platypus import Table, TableStyle
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                elements.append(t)
                elements.append(Spacer(1, 0.3))
        
        # Build PDF
        pdf_doc.build(elements)
        
        # Clean up
        pdf_buffer.seek(0)
        os.remove(docx_path)
        
        # Return the PDF
        user_id = session.get('user_id')
        log_progress(user_id, 'Converted Word to PDF')
        return send_file(
            pdf_buffer,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=secure_filename(file.filename.rsplit('.', 1)[0] + '.pdf')
        )
    
    except ImportError as e:
        return jsonify({'error': 'Required library not installed. Please install python-docx.'}), 500
    except Exception as e:
        print(f"Error converting Word to PDF: {str(e)}")
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500


@app.route('/convert-ppt-to-pdf', methods=['POST'])
def convert_ppt_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error':'No file provided'}),400
    file = request.files['file']
    if file.filename == '' or not file.filename.lower().endswith(('.pptx','.ppt')):
        return jsonify({'error':'Invalid file type'}),400
    try:
        from pptx import Presentation
        prs = Presentation(file)
        # simple text extraction to pdf
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        elems=[]
        for i,slide in enumerate(prs.slides, start=1):
            elems.append(Paragraph(f"Slide {i}", styles['Heading2']))
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    elems.append(Paragraph(shape.text, styles['Normal']))
            elems.append(Spacer(1,12))
        doc.build(elems)
        buffer.seek(0)
        user_id = session.get('user_id')
        log_progress(user_id, 'Converted PPT to PDF')
        return send_file(buffer, as_attachment=True, download_name=file.filename.rsplit('.',1)[0]+'.pdf', mimetype='application/pdf')
    except ImportError:
        return jsonify({'error':'python-pptx not installed'}),500
    except Exception as e:
        print('PPT->PDF error', e)
        return jsonify({'error':str(e)}),500


if __name__ == '__main__':
    app.run(debug=True)